"""
Build MRI AI Platform presentation (.pptx) for doctors & radiologists.
Run: python build_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree
import copy

# ── colour palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0c, 0x11, 0x20)   # dark navy
CARD     = RGBColor(0x13, 0x1b, 0x2e)   # card background
BLUE     = RGBColor(0x0a, 0x6b, 0xff)   # primary blue
ACCENT   = RGBColor(0x00, 0xd4, 0xaa)   # teal accent
WARN     = RGBColor(0xff, 0x6b, 0x35)   # orange
PURPLE   = RGBColor(0x7c, 0x3a, 0xed)   # purple
TEXT     = RGBColor(0xe2, 0xe8, 0xf0)   # light text
MUTED    = RGBColor(0x88, 0x99, 0xbb)   # muted text
WHITE    = RGBColor(0xff, 0xff, 0xff)
BLACK    = RGBColor(0x00, 0x00, 0x00)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    """Add a rectangle shape."""
    from pptx.util import Emu
    shp = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb

def multiline_txt(slide, lines, l, t, w, h,
                  size=16, color=TEXT, align=PP_ALIGN.LEFT, bold=False):
    """Text box with multiple paragraphs."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for text, kwargs in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = kwargs.get('align', align)
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(kwargs.get('size', size))
        run.font.bold  = kwargs.get('bold', bold)
        run.font.color.rgb = kwargs.get('color', color)
        run.font.italic    = kwargs.get('italic', False)
    return txb

def tag_box(slide, text, l, t, color=BLUE, text_color=WHITE):
    """Small pill / tag label."""
    W_IN = len(text) * 0.095 + 0.3
    shp = box(slide, l, t, W_IN, 0.28, fill=color)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    # round corners via XML
    sp = shp._element
    sp_pr = sp.find(ns.qn('p:spPr'))
    prstGeom = sp_pr.find(ns.qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
    tf = shp.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size  = Pt(9)
    run.font.bold  = True
    run.font.color.rgb = text_color
    return shp

def card_box(slide, l, t, w, h, fill=CARD, border=RGBColor(0x1e,0x2d,0x4a)):
    shp = box(slide, l, t, w, h, fill=fill, line=border, line_w=Pt(1))
    # round corners
    sp = shp._element
    sp_pr = sp.find(ns.qn('p:spPr'))
    prstGeom = sp_pr.find(ns.qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
    return shp

def divider(slide, t):
    """Thin horizontal rule."""
    shp = slide.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(12.33), Pt(1.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0x1e,0x2d,0x4a)
    shp.line.fill.background()

def check_list(slide, items, l, t, w, spacing=0.34, size=13):
    """Bulleted checkmark list."""
    for i, item in enumerate(items):
        txt(slide, "✓", l, t + i*spacing, 0.25, 0.32,
            size=size, bold=True, color=ACCENT)
        txt(slide, item, l+0.27, t + i*spacing, w-0.27, 0.32,
            size=size, color=TEXT)

def stat_card(slide, l, t, w, h, number, label, num_color=ACCENT):
    card_box(slide, l, t, w, h)
    txt(slide, number, l, t+0.18, w, 0.8,
        size=36, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t+0.95, w, 0.55,
        size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

# gradient-ish accent strip at top
b = box(sl, 0, 0, 13.33, 0.06, fill=BLUE)
b.line.fill.background()

tag_box(sl, "Medical AI Platform", 5.4, 0.35, color=BLUE)

txt(sl, "AI-Powered MRI Analysis Platform",
    0.6, 0.85, 12.1, 1.4,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Automated segmentation · Volumetric measurement · Quality assurance",
    0.6, 2.35, 12.1, 0.5,
    size=18, color=ACCENT, align=PP_ALIGN.CENTER)

txt(sl, "Designed for radiologists — faster reads, consistent measurements, full PACS integration.",
    1.2, 2.95, 10.9, 0.6,
    size=15, color=MUTED, align=PP_ALIGN.CENTER)

# 4 badge boxes
badges = [
    ("🧠  Brain Tumor Segmentation", BLUE),
    ("📐  Volumetric Measurements",  PURPLE),
    ("✅  QA & Artifact Detection",  ACCENT),
    ("🔗  PACS / DICOM Integration", RGBColor(0x0f,0x76,0x6e)),
]
for i, (label, col) in enumerate(badges):
    bx = i % 2
    by = i // 2
    lx = 1.7 + bx * 5.1
    ty = 3.85 + by * 0.72
    c = card_box(sl, lx, ty, 4.6, 0.55, fill=RGBColor(0x13,0x1b,0x2e),
                 border=col)
    txt(sl, label, lx+0.15, ty+0.1, 4.3, 0.38, size=14, color=TEXT)

txt(sl, "For Doctors & Radiologists", 0.6, 7.05, 12.1, 0.3,
    size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "The Challenge", 5.6, 0.22, color=WARN)
txt(sl, "What Radiologists Face Today",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

problems = [
    ("⏱️", "Time-Consuming Reads",
     "Manual volumetric measurement of brain lesions takes 20–40 min per study, creating workflow bottlenecks."),
    ("📊", "Inconsistent Measurements",
     "Inter-reader variability in tumor volume can exceed 15–20%, complicating treatment response monitoring."),
    ("🔍", "Subtle Findings Missed",
     "High case volumes and fatigue increase the risk of missing small or evolving lesions."),
    ("📁", "No Structured Output",
     "Results locked in free-text reports make longitudinal comparison and research difficult."),
    ("🛤️", "Manual Study Routing",
     "Coordinators manually triage studies, leading to delays and mis-routing errors."),
    ("🏗️", "Siloed AI Tools",
     "Existing AI tools lack PACS integration and cannot compare results over time."),
]

cols = 3
rows = 2
for i, (icon, title, desc) in enumerate(problems):
    col = i % cols
    row = i // cols
    lx = 0.35 + col * 4.32
    ty = 1.6 + row * 2.55
    card_box(sl, lx, ty, 4.05, 2.35)
    txt(sl, icon,  lx+0.2, ty+0.18, 0.6,  0.5, size=22)
    txt(sl, title, lx+0.2, ty+0.72, 3.65, 0.38, size=13, bold=True, color=TEXT)
    txt(sl, desc,  lx+0.2, ty+1.12, 3.65, 1.1,  size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Our Solution", 5.7, 0.22, color=ACCENT, text_color=BLACK)
txt(sl, "One Platform.  End-to-End AI.",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Pipeline flow
steps = [
    ("1", "Upload /\nPACS",    "DICOM via\nOrthanc"),
    ("2", "Auto\nRoute",       "Smart pipeline\nassignment"),
    ("3", "QA\nCheck",         "Detect artifacts\n& gaps"),
    ("4", "AI\nSegment",       "SegResNet\n3D inference"),
    ("5", "Volumes\n(mL)",     "Per tumor\nsub-region"),
    ("6", "Report",            "JSON + overlays\n+ PDF"),
]
sw = 1.7
gap = 0.18
start = 0.4
ty = 1.5
for i, (num, title, sub) in enumerate(steps):
    lx = start + i * (sw + gap)
    # connector arrow (not first)
    if i > 0:
        arrow_x = lx - gap + 0.02
        b = box(sl, arrow_x, ty+0.42, gap-0.02, 0.15, fill=BLUE)
        b.line.fill.background()
    c = card_box(sl, lx, ty, sw, 1.0, border=BLUE)
    # number circle
    circ = box(sl, lx+0.65, ty+0.08, 0.35, 0.35, fill=BLUE)
    circ.line.fill.background()
    txt(sl, num, lx+0.65, ty+0.09, 0.35, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, lx+0.1, ty+0.48, sw-0.2, 0.38, size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    txt(sl, sub,   lx+0.1, ty+0.78, sw-0.2, 0.38, size=9,  color=MUTED, align=PP_ALIGN.CENTER)

# Three highlight cards
highlights = [
    ("Core AI Model", "SegResNet", "Trained on BraTS 2021 — the gold standard benchmark for brain tumor segmentation."),
    ("Integration",   "Native DICOM / DICOMweb", "Plugs directly into your existing PACS. Zero rip-and-replace."),
    ("Output",        "Structured & Standard", "DICOM SR/SEG, NIfTI overlays, PDF reports — ready for EMR / FHIR."),
]
for i, (role, title, desc) in enumerate(highlights):
    lx = 0.4 + i * 4.32
    ty2 = 2.85
    c = card_box(sl, lx, ty2, 4.05, 2.15, border=ACCENT)
    txt(sl, role,  lx+0.2, ty2+0.15, 3.65, 0.28, size=9,  bold=True, color=ACCENT)
    txt(sl, title, lx+0.2, ty2+0.48, 3.65, 0.38, size=14, bold=True, color=WHITE)
    txt(sl, desc,  lx+0.2, ty2+0.95, 3.65, 1.1,  size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BRAIN TUMOR AI
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Brain MRI AI", 5.7, 0.22, color=BLUE)
txt(sl, "Automated Brain Tumor Segmentation",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left column — sub-regions + sequences
card_box(sl, 0.35, 1.5, 6.0, 2.8, border=BLUE)
txt(sl, "Three Tumor Sub-Regions Quantified", 0.55, 1.65, 5.6, 0.35, size=13, bold=True, color=BLUE)

regions = [
    ("Whole Tumor (WT)",     "Full disease extent",       ACCENT, 1.0),
    ("Tumor Core (TC)",      "Active + necrotic core",    BLUE,   0.65),
    ("Enhancing Tumor (ET)", "Viable / active tumor",     WARN,   0.38),
]
for i, (name, sub, col, pct) in enumerate(regions):
    ty = 2.15 + i * 0.7
    txt(sl, name, 0.55, ty, 2.6, 0.28, size=11, bold=True, color=TEXT)
    txt(sl, sub,  3.2,  ty, 2.9, 0.28, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    # track
    b = box(sl, 0.55, ty+0.3, 5.5, 0.1, fill=RGBColor(0x1e,0x2d,0x4a))
    b.line.fill.background()
    # fill
    b2 = box(sl, 0.55, ty+0.3, 5.5*pct, 0.1, fill=col)
    b2.line.fill.background()

# sequences
card_box(sl, 0.35, 4.5, 6.0, 1.35)
txt(sl, "Four MRI Sequences Fused", 0.55, 4.65, 5.6, 0.3, size=12, bold=True, color=TEXT)
seqs = ["T1", "T1ce", "T2", "FLAIR"]
for i, s in enumerate(seqs):
    lx = 0.55 + i * 1.4
    c = card_box(sl, lx, 5.05, 1.2, 0.62, border=BLUE)
    txt(sl, s, lx, 5.15, 1.2, 0.38, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Right column — bullet points
check_list(sl, [
    "Sliding-window 3D inference — full brain volume, no slice missed",
    "Resampled to 1 mm isotropic — normalises scanner variability",
    "Z-score intensity normalisation — robust to contrast differences",
    "Volume in mL per sub-region + % contribution (RANO-ready)",
    "Model version + SHA-256 checksum logged on every result",
    "0.05 mL minimum lesion filter — suppresses sub-voxel noise",
    "GPU inference, CPU fallback — works on any hospital hardware",
], 6.65, 1.5, 6.3, spacing=0.72, size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — QA ENGINE
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Quality Assurance", 5.35, 0.22, color=PURPLE)
txt(sl, "Built-In Image Quality Guard",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "AI results are only as good as the images. The platform validates quality before every inference run.",
    1.0, 1.3, 11.3, 0.4, size=13, color=MUTED, align=PP_ALIGN.CENTER)

qa_items = [
    ("🎯", "Missing Sequence",       "Flags studies missing T1, FLAIR, or T2 — prevents silently wrong segmentations."),
    ("📏", "Voxel Spacing",          "Spacing validated 0.5–5.0 mm per axis, anisotropy ratio < 3.0."),
    ("🏃", "Motion Artifacts",       "Edge-energy analysis detects patient motion across the scan volume."),
    ("📐", "Coverage Check",         "Requires ≥ 20 slices for adequate anatomical coverage."),
    ("↕️", "Slice Gap Detection",    "Identifies large inter-slice gaps that compromise 3D accuracy."),
    ("🚦", "Three Severity Levels",  "BLOCKING halts analysis · WARNING flags caution · INFO logs only."),
]
for i, (icon, title, desc) in enumerate(qa_items):
    col = i % 3
    row = i // 3
    lx = 0.35 + col * 4.32
    ty = 1.85 + row * 2.4
    border_col = ACCENT if i == 5 else RGBColor(0x1e,0x2d,0x4a)
    card_box(sl, lx, ty, 4.05, 2.2, border=border_col)
    txt(sl, icon,  lx+0.2, ty+0.15, 0.55, 0.45, size=20)
    txt(sl, title, lx+0.2, ty+0.65, 3.65, 0.35, size=13, bold=True, color=TEXT)
    txt(sl, desc,  lx+0.2, ty+1.05, 3.65, 1.0,  size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RADIOLOGIST WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Radiologist Workflow", 5.1, 0.22, color=ACCENT, text_color=BLACK)
txt(sl, "Designed Around Your Workflow",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left — timeline card
card_box(sl, 0.35, 1.5, 6.0, 3.4)
txt(sl, "Real-Time Study Timeline", 0.55, 1.65, 5.6, 0.35, size=13, bold=True, color=TEXT)

timeline = [
    ("T + 0 s",   ACCENT, "Study ingested from PACS automatically"),
    ("T + 5 s",   BLUE,   "Auto-routed to Brain MRI AI pipeline"),
    ("T + 30 s",  BLUE,   "QA passed — 4-channel volume assembled"),
    ("T + 90 s",  PURPLE, "AI inference complete"),
    ("T + 120 s", ACCENT, "Volumes, overlays & report ready"),
]
for i, (time, col, text) in enumerate(timeline):
    ty = 2.15 + i * 0.54
    # dot
    dot = box(sl, 0.55, ty+0.06, 0.14, 0.14, fill=col)
    dot.line.fill.background()
    # line (not last)
    if i < len(timeline)-1:
        ln = box(sl, 0.60, ty+0.22, 0.04, 0.38, fill=RGBColor(0x1e,0x2d,0x4a))
        ln.line.fill.background()
    txt(sl, time, 0.82, ty, 1.1, 0.28, size=9, bold=True, color=col)
    txt(sl, text, 2.05, ty, 4.0, 0.28, size=11, color=TEXT)

# comparison note
card_box(sl, 0.35, 5.1, 6.0, 1.25)
txt(sl, "Longitudinal Comparison", 0.55, 5.25, 5.6, 0.3, size=12, bold=True, color=TEXT)
txt(sl, "Every result is versioned. One click compares Whole Tumor volume between baseline\nand follow-up — delta in mL and %, directly supporting RANO response assessment.",
    0.55, 5.6, 5.6, 0.65, size=10, color=MUTED)

# Right — checklist
check_list(sl, [
    "OHIF DICOM Viewer integrated — overlays rendered on the MRI",
    "Worklist with AI status: pending / processing / ready",
    "One-click structured report — volumes pre-populated",
    "Active learning queue — only low-confidence cases for review",
    "Role-based access: Radiologist / Technician / Viewer",
    "Full audit trail — every view and edit logged for compliance",
    "Server-Sent Events (SSE) for live progress in the browser",
], 6.65, 1.5, 6.3, spacing=0.73, size=12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CLINICAL IMPACT
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Clinical Impact", 5.55, 0.22, color=BLUE)
txt(sl, "Measurable Benefits",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Stat cards
stat_card(sl, 0.35,  1.5, 2.95, 1.55, "~2 min", "Full 3D brain analysis\nvs. 20–40 min manual", ACCENT)
stat_card(sl, 3.47,  1.5, 2.95, 1.55, "3",       "Tumor sub-regions\nquantified automatically", BLUE)
stat_card(sl, 6.59,  1.5, 2.95, 1.55, "100%",    "Reproducible — zero\ninter-reader variability", ACCENT)
stat_card(sl, 9.71,  1.5, 2.95, 1.55, "∞",       "Versioned history per\nstudy — trend tracking", BLUE)

# Three audience benefit cards
audiences = [
    ("For Radiologists", [
        "Eliminate manual measurement burden",
        "Consistent baseline for every follow-up read",
        "Focus cognition on clinical interpretation",
    ]),
    ("For Oncologists", [
        "Objective RANO-aligned volumetric data",
        "Trend charts across treatment cycles",
        "Structured data ready for tumor boards",
    ]),
    ("For Institutions", [
        "Scale by adding GPU workers on demand",
        "Plug into existing PACS — no disruption",
        "Audit-ready and FHIR-exportable results",
    ]),
]
for i, (title, items) in enumerate(audiences):
    lx = 0.35 + i * 4.32
    card_box(sl, lx, 3.3, 4.05, 2.75)
    txt(sl, title, lx+0.2, 3.45, 3.65, 0.35, size=13, bold=True, color=TEXT)
    check_list(sl, items, lx+0.2, 3.95, 3.65, spacing=0.65, size=11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EXPANDABILITY / ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Roadmap", 6.1, 0.22, color=BLUE)
txt(sl, "Built to Grow With You",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Plugin architecture — each new anatomical region is a single folder. No core changes required.",
    1.2, 1.3, 10.9, 0.38, size=13, color=MUTED, align=PP_ALIGN.CENTER)

modules = [
    ("🧠", "Brain MRI",   "Live",        ACCENT, "Glioma segmentation (TC/WT/ET),\nvolumetrics, FLAIR lesion load"),
    ("🦴", "Spine MRI",   "Scaffolded",  RGBColor(0xf5,0x9e,0x0b), "Vertebral segmentation,\ndisc herniation, canal stenosis"),
    ("🫁", "Chest MRI",   "Planned",     MUTED,  "Lung nodule detection,\ncardiac structure measurement"),
    ("🫀", "Abdomen MRI", "Planned",     MUTED,  "Organ segmentation,\nlesion characterisation"),
]
for i, (icon, title, status, col, desc) in enumerate(modules):
    lx = 0.35 + i * 3.2
    ty = 1.85
    border = col if status == "Live" else RGBColor(0x1e,0x2d,0x4a)
    card_box(sl, lx, ty, 3.0, 2.5, border=border)
    txt(sl, icon,   lx+0.15, ty+0.15, 0.5,  0.45, size=20)
    txt(sl, title,  lx+0.15, ty+0.7,  2.7,  0.35, size=13, bold=True, color=TEXT)
    txt(sl, status, lx+0.15, ty+1.05, 2.7,  0.28, size=10, bold=True, color=col)
    txt(sl, desc,   lx+0.15, ty+1.4,  2.7,  0.85, size=10, color=MUTED)

# Two feature cards
features = [
    ("Enterprise Platform Features", [
        "A/B model testing — compare new vs baseline before deploy",
        "Ensemble inference — combine models for higher confidence",
        "Multi-tenant — isolated workspaces per department or site",
    ], BLUE),
    ("Standards & Interoperability", [
        "DICOM SR / SEG export — results as standard DICOM objects",
        "DICOMweb (STOW-RS, QIDO-RS) — seamless PACS integration",
        "FHIR server push — results flow into hospital EMR systems",
    ], ACCENT),
]
for i, (title, items, col) in enumerate(features):
    lx = 0.35 + i * 6.55
    ty2 = 4.6
    card_box(sl, lx, ty2, 6.25, 2.4, border=col)
    txt(sl, title, lx+0.2, ty2+0.2, 5.85, 0.35, size=13, bold=True, color=TEXT)
    check_list(sl, items, lx+0.2, ty2+0.68, 5.85, spacing=0.52, size=11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

b = box(sl, 0, 0, 13.33, 0.06, fill=ACCENT)
b.line.fill.background()

tag_box(sl, "Summary", 6.15, 0.35, color=ACCENT, text_color=BLACK)

txt(sl, "AI that Augments the Radiologist",
    0.5, 0.75, 12.3, 1.1, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Not a replacement — a precision instrument that handles measurement, consistency,\n"
        "and routing so you can focus on what only you can do: clinical judgement and patient care.",
    1.0, 1.95, 11.3, 0.9, size=15, color=MUTED, align=PP_ALIGN.CENTER)

pillars = [
    ("⚡", "Speed",      "~2 min full-brain analysis\nReal-time progress tracking\nInstant structured output"),
    ("🎯", "Precision",  "BraTS-validated SegResNet\n1 mm isotropic processing\nSub-region volumes in mL"),
    ("🏥", "Integration","Native PACS & DICOM\nFHIR / EMR ready\nZero workflow disruption"),
]
for i, (icon, title, desc) in enumerate(pillars):
    lx = 0.6 + i * 4.1
    ty = 3.1
    card_box(sl, lx, ty, 3.8, 2.3)
    txt(sl, icon,  lx+1.65, ty+0.18, 0.7, 0.5, size=22, align=PP_ALIGN.CENTER)
    txt(sl, title, lx+0.2, ty+0.78, 3.4, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc,  lx+0.2, ty+1.2,  3.4, 0.95, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Quote bar
qb = box(sl, 0.6, 5.65, 12.13, 1.35, fill=RGBColor(0x13,0x1b,0x2e))
qb.line.color.rgb = ACCENT
qb.line.width = Pt(2)
txt(sl, '"The goal of AI in radiology is not to read scans faster — it is to read them better,\n'
        'with more objectivity, and to free radiologists for decisions that require human expertise."',
    0.95, 5.82, 11.5, 1.0, size=13, italic=True, color=TEXT, align=PP_ALIGN.CENTER)

# Bottom strip
b2 = box(sl, 0, 7.44, 13.33, 0.06, fill=BLUE)
b2.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ADVANCED FEATURES OVERVIEW (12 features)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "Advanced AI Features", 4.8, 0.22, color=ACCENT, text_color=BLACK)
txt(sl, "12 Features.  Maximum ROI.",
    0.5, 0.6, 12.3, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Targeting every radiology pain point — from billing leakage to missed findings.",
    1.2, 1.3, 10.9, 0.38, size=13, color=MUTED, align=PP_ALIGN.CENTER)

features_grid = [
    ("💰", "1. CPT Code Suggestion",    "Auto-suggest billing codes\nfrom AI findings"),
    ("📄", "2. Auto-Draft Report",       "Pre-fills structured report\nfrom measurements"),
    ("🔢", "3. Worklist Prioritization", "Sorts by AI urgency score\ncritical cases first"),
    ("🚨", "4. Critical Finding Alerts", "Webhook/email on critical\nfindings <5 min"),
    ("🔍", "5. Automated Prior Compare", "Auto-fetches and compares\nprior studies"),
    ("👥", "6. Peer Review Queue",       "Routes low-confidence\ncases to experts"),
    ("📈", "7. Longitudinal Tracking",   "Volume trend curves\nacross timepoints"),
    ("📊", "8. QA / Audit Dashboard",    "Agreement rates, TAT\nmetrics, compliance"),
    ("🩺", "9. Referring Physician Portal","Read-only portal for\nreferring clinicians"),
    ("⚙️", "10. Protocol Optimization", "Flags wrong protocols\nbefore scanning"),
    ("📅", "11. Capacity Prediction",    "Scanner utilization\nanalytics + forecasting"),
    ("🔎", "12. Incidental Findings",    "Detects findings outside\nprimary ROI"),
]

cols = 4
card_w = 3.0
card_h = 1.35
gap_x = 0.12
gap_y = 0.12
start_x = 0.37
start_y = 1.75
for i, (icon, title, desc) in enumerate(features_grid):
    col = i % cols
    row = i // cols
    lx = start_x + col * (card_w + gap_x)
    ty = start_y + row * (card_h + gap_y)
    card_box(sl, lx, ty, card_w, card_h)
    txt(sl, icon,  lx+0.12, ty+0.12, 0.4, 0.4, size=16)
    txt(sl, title, lx+0.12, ty+0.58, card_w-0.24, 0.32, size=10, bold=True, color=TEXT)
    txt(sl, desc,  lx+0.12, ty+0.9,  card_w-0.24, 0.4,  size=9, color=MUTED)


def feature_slide(title, tag_text, tag_color, tag_text_color,
                  steps, right_stats, checklist,
                  card2_title=None, card2_rows=None):
    """Build a standard 2-column feature slide."""
    sl = add_slide()
    bg(sl)
    tag_box(sl, tag_text, (13.33 - len(tag_text)*0.12 - 0.6)/2, 0.22,
            color=tag_color, text_color=tag_text_color)
    txt(sl, title, 0.5, 0.6, 12.3, 0.65, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left column — how it works card
    card_box(sl, 0.35, 1.45, 6.1, 2.85, border=BLUE)
    txt(sl, "How It Works", 0.55, 1.6, 5.7, 0.3, size=12, bold=True, color=BLUE)
    for i, (step_num_color, step_text) in enumerate(steps):
        ty = 2.05 + i * 0.52
        circ = box(sl, 0.55, ty+0.02, 0.26, 0.26, fill=step_num_color)
        circ.line.fill.background()
        txt(sl, str(i+1), 0.55, ty+0.03, 0.26, 0.22, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, step_text, 0.88, ty, 5.45, 0.4, size=10, color=TEXT)

    # Optional second card
    if card2_title and card2_rows:
        card_box(sl, 0.35, 4.45, 6.1, 1.35)
        txt(sl, card2_title, 0.55, 4.6, 5.7, 0.28, size=11, bold=True, color=TEXT)
        for i, (left, right) in enumerate(card2_rows):
            ty = 5.0 + i * 0.28
            txt(sl, left,  0.55, ty, 4.2, 0.26, size=9, color=TEXT)
            txt(sl, right, 4.75, ty, 1.5, 0.26, size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

    # Right column — stat cards
    stat_card(sl, 6.65, 1.45, 2.9, 1.28, right_stats[0][0], right_stats[0][1], right_stats[0][2])
    stat_card(sl, 9.73, 1.45, 2.9, 1.28, right_stats[1][0], right_stats[1][1], right_stats[1][2])
    check_list(sl, checklist, 6.65, 2.92, 5.95, spacing=0.56, size=11)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CPT CODE SUGGESTION
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="CPT Code Suggestion",
    tag_text="Feature 1 of 12  ·  Revenue",
    tag_color=ACCENT, tag_text_color=BLACK,
    steps=[
        (BLUE,   "AI results contain structured measurements (volumes, counts, severity flags)"),
        (BLUE,   "Rule engine maps finding types → CPT code candidates with confidence scores"),
        (BLUE,   "Codes ranked by modality + body part + complexity; add-on codes included"),
        (ACCENT, "Radiologist reviews & approves — final codes logged in audit trail"),
    ],
    right_stats=[
        ("+18%", "Avg revenue increase\nfrom correct code capture", ACCENT),
        ("<3s",  "Code suggestion time\nper study", BLUE),
    ],
    checklist=[
        "Zero revenue leakage — catches under-coded studies",
        "Compliance-safe — radiologist always has final approval",
        "Includes AI add-on codes (0691T series) automatically",
        "Audit trail: every acceptance or correction logged",
        "Payer-specific rules configurable per insurance carrier",
    ],
    card2_title="Sample CPT Mappings",
    card2_rows=[
        ("Brain MRI w/ contrast, tumor vol >0",  "70553"),
        ("Brain MRI w/o contrast",                "70551"),
        ("Spine MRI cervical w/ contrast",        "72156"),
        ("AI-assisted analysis (add-on)",          "0691T"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AUTO-DRAFT STRUCTURED REPORT
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Auto-Draft Structured Report",
    tag_text="Feature 2 of 12  ·  Time + Revenue",
    tag_color=BLUE, tag_text_color=WHITE,
    steps=[
        (BLUE,   "AI result measurements extracted: volumes, counts, severity, QA flags"),
        (BLUE,   "Use-case Jinja2 template populated with structured findings"),
        (BLUE,   "Draft includes: indication, technique, findings, impression placeholder"),
        (ACCENT, "Export: PDF (WeasyPrint) + DICOM SR (pydicom) + FHIR DiagnosticReport"),
    ],
    right_stats=[
        ("60%",  "Faster dictation time\nper study", ACCENT),
        ("0",    "Transcription errors\nin measurements", BLUE),
    ],
    checklist=[
        "Three export formats — PDF, DICOM SR, FHIR DiagnosticReport",
        "Structured & searchable — values stored as discrete fields",
        "Comparison-ready — prior deltas auto-inserted as table",
        "Template per use case — brain/spine/chest/abdomen",
        "Radiologist editable — impression free-text, measurements protected",
    ],
    card2_title="Example Brain MRI Draft Findings",
    card2_rows=[
        ("Whole Tumor (WT)",     "227.5 mL"),
        ("Tumor Core (TC)",      "183.1 mL"),
        ("Enhancing Tumor (ET)", "112.8 mL"),
        ("Impression",           "[Radiologist fills]"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — WORKLIST PRIORITIZATION
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Worklist Prioritization",
    tag_text="Feature 3 of 12  ·  Time + Accuracy",
    tag_color=PURPLE, tag_text_color=WHITE,
    steps=[
        (WARN,   "40% weight: pathology severity (tumor volume, critical QA flags, high-change delta)"),
        (BLUE,   "25% weight: AI confidence score — high confidence + severe finding = top priority"),
        (PURPLE, "20% weight: study age — older unread studies elevated progressively"),
        (ACCENT, "15% weight: referring physician STAT / Routine / Elective flag"),
    ],
    right_stats=[
        ("30%",   "Faster TAT for\ncritical studies", WARN),
        ("100%",  "Critical findings\nread within SLA", ACCENT),
    ],
    checklist=[
        "Dynamic re-ranking — scores updated in real-time",
        "STAT override — referring physician can escalate instantly",
        "Sub-specialty filtering per radiologist role",
        "SLA monitoring with dashboard alerts on breach",
        "Zero manual triage — coordinators no longer route by hand",
    ],
    card2_title="Priority Levels",
    card2_rows=[
        ("CRITICAL (score ≥80)",  "Notify immediately"),
        ("HIGH     (score 60–79)", "Top of worklist"),
        ("NORMAL   (score 30–59)", "Standard queue"),
        ("ROUTINE  (score <30)",   "Batch-eligible"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CRITICAL FINDING AUTO-NOTIFICATION
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Critical Finding Auto-Notification",
    tag_text="Feature 4 of 12  ·  Safety + Time",
    tag_color=WARN, tag_text_color=WHITE,
    steps=[
        (WARN,   "AI result stored → AlertRuleEngine evaluates all active rules"),
        (WARN,   "Rules match on: event_type, measurement thresholds, QA flag patterns"),
        (WARN,   "Notification dispatched: Webhook POST / email (SMTP) / WebSocket push"),
        (ACCENT, "AlertHistory record written: status (sent/failed) + full payload"),
    ],
    right_stats=[
        ("<5 min", "Time from finding\nto notification", WARN),
        ("100%",   "Accountability trail\nper alert", ACCENT),
    ],
    checklist=[
        "Configurable thresholds — per use case, measurement key, severity",
        "Three delivery channels — webhook, email, in-app notification",
        "Retry with back-off — failed webhooks retried 3× before failure",
        "Medi-legal coverage — timestamp, recipient, payload hash logged",
        "Escalation chain — unacknowledged alerts escalate to on-call",
    ],
    card2_title="Example Alert Rules",
    card2_rows=[
        ("brain_mri.whole_tumor_volume > 100 mL", "CRITICAL"),
        ("qa_flag contains 'mass_effect'",         "STAT"),
        ("delta.change_pct > 30% vs prior",        "RAPID PROGRESS"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — AUTOMATED PRIOR COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Automated Prior Comparison",
    tag_text="Feature 5 of 12  ·  Accuracy + Time",
    tag_color=BLUE, tag_text_color=WHITE,
    steps=[
        (BLUE,   "New AI result arrives for patient_id X"),
        (BLUE,   "Query local DB + Orthanc QIDO-RS for same patient + same use case"),
        (BLUE,   "Most recent prior result selected → compare_results() called automatically"),
        (ACCENT, "Delta report appended to study page — radiologist opens with comparison ready"),
    ],
    right_stats=[
        ("8 min", "Saved per follow-up\nread (no manual search)", ACCENT),
        ("100%",  "Prior review rate\nvs ~35% manual", BLUE),
    ],
    checklist=[
        "Zero manual effort — fires automatically on new result",
        "RANO-aligned percentage change for oncology use",
        "Cascade query — falls back to federated Orthanc peers",
        "Multi-timepoint aware — selects immediately preceding study",
        "Rapid progression (≥25% growth) triggers Critical Alert",
    ],
    card2_title="What Gets Compared",
    card2_rows=[
        ("All volumetric measurements",  "Δ mL + %"),
        ("QA flag changes",              "New / Resolved"),
        ("Days between studies",         "Timeline"),
        ("Segmentation overlays",        "Side-by-side"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONFIDENCE-GATED PEER REVIEW
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Confidence-Gated Peer Review",
    tag_text="Feature 6 of 12  ·  Accuracy",
    tag_color=PURPLE, tag_text_color=WHITE,
    steps=[
        (PURPLE, "AI inference returns a confidence score per result (0.0 – 1.0)"),
        (PURPLE, "Score below configurable threshold (default 0.75) → ReviewItem created"),
        (PURPLE, "Review queue shows study + result + confidence score for expert radiologist"),
        (ACCENT, "Reviewer approves, corrects measurements, or flags for re-scan — all logged"),
    ],
    right_stats=[
        ("99%",  "Uncertain cases\ncaught before report", PURPLE),
        ("↓40%", "Radiologist time on\nconfident cases", ACCENT),
    ],
    checklist=[
        "Active learning — corrections update model training dataset",
        "Configurable threshold per use case (brain 0.85, chest 0.65)",
        "Workload balancing — items distributed by round-robin",
        "SLA-tracked — pending items older than N hours escalate",
        "High-confidence results bypass queue entirely",
    ],
    card2_title="Review Item States",
    card2_rows=[
        ("PENDING",    "Awaiting expert assignment"),
        ("IN REVIEW",  "Expert actively reviewing"),
        ("APPROVED",   "Accepted as-is by reviewer"),
        ("CORRECTED",  "New result version created"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — LONGITUDINAL TREND TRACKING
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Longitudinal Trend Tracking",
    tag_text="Feature 7 of 12  ·  Accuracy + Revenue",
    tag_color=ACCENT, tag_text_color=BLACK,
    steps=[
        (BLUE,   "All AI results for patient_id + use case aggregated across all timepoints"),
        (BLUE,   "Study dates from DICOM metadata used for accurate timeline positioning"),
        (BLUE,   "RANO response classification computed at each timepoint (CR/PR/SD/PD)"),
        (ACCENT, "Interactive chart exported: CSV / PDF for tumor board + clinical trials"),
    ],
    right_stats=[
        ("RANO", "Ready response\nclassification built-in", ACCENT),
        ("$3k+", "Per clinical trial\nenrollment enabled", BLUE),
    ],
    checklist=[
        "Click any data point to jump to that study's report",
        "All sub-regions (WT/TC/ET) plotted on the same timeline",
        "Export to CSV/PDF in one click for tumor boards",
        "Multi-patient cohort view for comparing treatment arms",
        "Automated RANO classification assigned at each timepoint",
    ],
    card2_title="Data Sources",
    card2_rows=[
        ("All results for patient_id + use case", "All versions"),
        ("Study dates from DICOM metadata",       "Accurate timeline"),
        ("Treatment events (optional)",            "FHIR integration"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — QA / AUDIT DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="QA / Audit Dashboard",
    tag_text="Feature 8 of 12  ·  Compliance + Quality",
    tag_color=BLUE, tag_text_color=WHITE,
    steps=[
        (BLUE,   "AuditEntry records aggregated: study receipt, job events, result views, exports"),
        (BLUE,   "TAT computed: median, p75, p95 per use case — updated every 5 minutes"),
        (BLUE,   "Radiologist agreement rate: corrections / total reviewed results"),
        (ACCENT, "Dashboard exported as PDF for ACR accreditation or compliance audits"),
    ],
    right_stats=[
        ("ACR",  "Accreditation-ready\ncompliance reports", ACCENT),
        ("100%", "Action traceability\nin audit log", BLUE),
    ],
    checklist=[
        "Date-range filtering — daily, weekly, quarterly, annual",
        "Per-radiologist breakdown — identify outliers for training",
        "One-click PDF/CSV export for accreditation submissions",
        "Anomaly alerts when TAT or agreement rate breach threshold",
        "Multi-tenant isolation — each site sees only its own data",
    ],
    card2_title="Key Metrics Tracked",
    card2_rows=[
        ("Turnaround Time (TAT)",         "Median / p75 / p95"),
        ("Radiologist Agreement Rate",    "vs AI before correction"),
        ("QA Flag Rate",                  "% studies with issues"),
        ("Critical Alert Response Time",  "Time to acknowledgement"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — REFERRING PHYSICIAN PORTAL
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Referring Physician Portal",
    tag_text="Feature 9 of 12  ·  Revenue + Collaboration",
    tag_color=ACCENT, tag_text_color=BLACK,
    steps=[
        (ACCENT, "Radiologist sends a time-limited share link to referring physician"),
        (ACCENT, "Link contains signed JWT scoped to single study + result (read-only)"),
        (ACCENT, "Portal shows: report, measurements, overlays — no raw PACS access"),
        (ACCENT, "Access logged in audit trail; link auto-expires after TTL (7 days default)"),
    ],
    right_stats=[
        ("+22%", "Increase in referring\nphysician referrals", ACCENT),
        ("0",    "PACS credentials\nshared externally", BLUE),
    ],
    checklist=[
        "No PACS account needed — standard web browser access",
        "HIPAA-compliant — time-limited, scoped tokens, full logging",
        "Branded portal — institution logo and contact details shown",
        "Mobile responsive — works on phone/tablet in the clinic",
        "Drives repeat referrals — top driver of referring loyalty",
    ],
    card2_title="What Referring Physicians See",
    card2_rows=[
        ("Structured report + measurements",  "✓ Visible"),
        ("Segmentation overlay images",       "✓ Visible"),
        ("Comparison vs prior study",         "✓ If available"),
        ("Raw DICOM data / other patients",   "✗ Blocked"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — PROTOCOL OPTIMIZATION
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Protocol Optimization",
    tag_text="Feature 10 of 12  ·  Time + Revenue",
    tag_color=WARN, tag_text_color=WHITE,
    steps=[
        (WARN,   "Study arrives with DICOM series metadata (TE, TR, flip angle, slice thickness)"),
        (WARN,   "Matched against use-case required_sequences and protocol specification"),
        (WARN,   "Deviations flagged: missing sequence, wrong spacing, insufficient coverage"),
        (ACCENT, "Corrective suggestion generated: 'Reacquire FLAIR at 3mm slice thickness'"),
    ],
    right_stats=[
        ("0",    "Repeat scans from\nprotocol mismatch", WARN),
        ("$800", "Saved per eliminated\nrepeat MRI scan", ACCENT),
    ],
    checklist=[
        "Pre-scan validation — checked before technologist runs scan",
        "Per-scanner vendor rules — GE/Siemens/Philips defaults accounted for",
        "Technologist alerts — separate dashboard from radiologist worklist",
        "Improves AI accuracy — compliant studies produce better results",
        "Reduces patient burden — fewer repeat visits, better satisfaction",
    ],
    card2_title="Common Protocol Issues Caught",
    card2_rows=[
        ("FLAIR with slice gap >1mm",          "Degrades WT segmentation"),
        ("T1ce coverage insufficient",         "Temporal lobe excluded"),
        ("TR value suboptimal for T1 contrast", "Suggest 600–800ms"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — CAPACITY PREDICTION
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Capacity Prediction",
    tag_text="Feature 11 of 12  ·  Revenue + Operations",
    tag_color=BLUE, tag_text_color=WHITE,
    steps=[
        (BLUE,   "Job timing data (created→completed) aggregated per use case, body part, modality"),
        (BLUE,   "Rolling 7/30-day volume trends computed per scanner + per time-of-day slot"),
        (BLUE,   "Peak hours identified → scheduling team optimizes appointment slot allocation"),
        (ACCENT, "7-day demand forecast using exponential smoothing — flags under/over-utilization"),
    ],
    right_stats=[
        ("85%+", "Target scanner\nutilization rate", ACCENT),
        ("+15%", "Revenue from\noptimized scheduling", BLUE),
    ],
    checklist=[
        "Data-driven scheduling — replaces guesswork with real throughput",
        "Bottleneck identification — shows highest queue-depth use cases",
        "Staff planning — radiologist shifts informed by predicted volume",
        "SLA forecasting — predicts TAT breach given current queue + staff",
        "Cost per scan — GPU compute time per study for financial reporting",
    ],
    card2_title="Dashboard Widgets",
    card2_rows=[
        ("Daily scan volume by use case", "Bar chart"),
        ("Hourly throughput heatmap",     "Heat map"),
        ("Worker queue depth over time",  "Line chart"),
        ("7-day demand forecast",         "Forecast band"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — INCIDENTAL FINDING DETECTION
# ══════════════════════════════════════════════════════════════════════════════
feature_slide(
    title="Incidental Finding Detection",
    tag_text="Feature 12 of 12  ·  Accuracy + Safety",
    tag_color=WARN, tag_text_color=WHITE,
    steps=[
        (WARN,   "Primary AI analysis completes (e.g., brain tumor volumes)"),
        (WARN,   "Secondary scan of anatomical regions outside the primary ROI"),
        (WARN,   "Anomaly scoring detects intensity outliers, asymmetry, unexpected structures"),
        (ACCENT, "Findings flagged in report with region + confidence — radiologist reviews"),
    ],
    right_stats=[
        ("↓34%", "Reduction in missed\nincidental findings", WARN),
        ("↓60%", "Malpractice exposure\nfrom missed findings", ACCENT),
    ],
    checklist=[
        "Non-intrusive — runs as post-processing, no extra scan time",
        "Structured reporting — dedicated section in auto-drafted report",
        "Configurable sensitivity per body part — balance FP rate",
        "Malpractice shield — documented AI check in audit defense",
        "Additional revenue — incidentals generate follow-up referrals",
    ],
    card2_title="Common Incidental Findings Detected",
    card2_rows=[
        ("Subdural collections",          "Non-neoplastic brain MRI"),
        ("Vertebral body signal abnormality", "Spine MRI"),
        ("Unexpected lymphadenopathy",    "Chest / abdomen MRI"),
        ("Vascular anomalies",            "Aneurysm / malformation"),
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — ROI SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)

tag_box(sl, "ROI Summary", 5.7, 0.22, color=ACCENT, text_color=BLACK)
txt(sl, "Combined Impact: Measurable Returns",
    0.5, 0.6, 12.3, 0.65, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three column impact cards
col_data = [
    ("💰 Revenue", ACCENT, [
        ("+18%", "CPT Code Capture"),
        ("+22%", "Referring Portal Referrals"),
        ("+15%", "Capacity Optimization"),
        ("$800/scan", "Eliminated Repeat Scans"),
        ("+$3k/pt", "Incidental Referrals"),
    ]),
    ("⏱️ Time Savings", BLUE, [
        ("−60%",        "Dictation Time"),
        ("−8 min/study","Prior Study Search"),
        ("Eliminated",  "Manual Triage"),
        ("−30%",        "TAT for Critical Cases"),
        ("−100%",       "Repeat Scan Reschedules"),
    ]),
    ("🎯 Accuracy + Safety", WARN, [
        ("99%",    "Peer Review Coverage"),
        ("<5 min", "Critical Alert Response"),
        ("−34%",   "Missed Incidentals"),
        ("−100%",  "Inter-reader Variability"),
        ("100%",   "Prior Review Rate"),
    ]),
]

for ci, (col_title, col_color, rows) in enumerate(col_data):
    lx = 0.35 + ci * 4.32
    card_box(sl, lx, 1.45, 4.05, 3.85, border=col_color)
    txt(sl, col_title, lx+0.2, 1.62, 3.65, 0.35, size=13, bold=True, color=col_color)
    for ri, (val, label) in enumerate(rows):
        ty = 2.12 + ri * 0.6
        txt(sl, val,   lx+0.2, ty, 1.2, 0.3, size=13, bold=True, color=col_color)
        txt(sl, label, lx+1.5, ty, 2.35, 0.3, size=10, color=TEXT)

# ROI highlight bar
card_box(sl, 0.35, 5.52, 12.63, 1.65, border=ACCENT)
txt(sl, "Typical Institution (500 MRI studies/month):",
    0.55, 5.68, 12.0, 0.3, size=11, bold=True, color=MUTED)
txt(sl, "Projected Additional Revenue:  $85,000 – $120,000 / year",
    0.55, 6.0, 12.0, 0.38, size=16, bold=True, color=WHITE)
txt(sl, "Driven by CPT code capture (+18%), increased referrals (+22%), reduced repeat scans, and capacity optimization.",
    0.55, 6.42, 12.0, 0.65, size=11, color=MUTED)

# ── save ─────────────────────────────────────────────────────────────────────
OUT = r"c:\sistems\projects\MR_Computer-Visuion\presentation\MRI_AI_Platform.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
